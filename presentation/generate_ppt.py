import sys
import os
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def main():
    # Load template
    template_path = "Idea Submission Template _ Redrob.pptx"
    if not os.path.exists(template_path):
        template_path = "../Idea Submission Template _ Redrob.pptx"
        
    prs = Presentation(template_path)
    
    # Read slides dump text
    dump_path = "scratch/slides_dump.txt"
    if not os.path.exists(dump_path):
        dump_path = "../scratch/slides_dump.txt"
        
    with open(dump_path, "r", encoding="utf-8") as f:
        content = f.read()
        
    # Parse slides using re.split
    parts = re.split(r"=== SLIDE (\d+) ===", content)
    slides_data = {}
    for i in range(1, len(parts), 2):
        slide_num = int(parts[i])
        slide_content = parts[i+1].strip()
        lines = [line.strip() for line in slide_content.split("\n") if line.strip()]
        slides_data[slide_num] = lines
        
    print(f"Parsed {len(slides_data)} slides from dump.")
    
    # Style definitions
    font_family = "Calibri"
    title_color = RGBColor(26, 54, 93)   # Dark Blue (#1A365D)
    sub_color = RGBColor(49, 151, 149)   # Teal (#319795)
    body_color = RGBColor(45, 55, 72)    # Charcoal (#2D3748)
    
    # Helper to style a paragraph
    def format_para(para, text, size, color, bold=False, italic=False, bullet=False):
        para.text = text
        para.font.name = font_family
        para.font.size = Pt(size)
        para.font.color.rgb = color
        para.font.bold = bold
        para.font.italic = italic
        if bullet:
            para.level = 1
        else:
            para.level = 0
            
    # Modify Slide 1
    s1 = prs.slides[0]
    s1_data = slides_data.get(1, [])
    team_name = ""
    problem_statement = ""
    leader_name = ""
    for line in s1_data:
        if line.startswith("Team Name:"):
            team_name = line
        elif line.startswith("Problem Statement:"):
            problem_statement = line
        elif line.startswith("Team Leader Name:"):
            leader_name = line
            
    for shape in s1.shapes:
        if shape.shape_id == 55 and team_name:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            format_para(p, team_name, 24, title_color, bold=True)
        elif shape.shape_id == 56 and problem_statement:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            format_para(p, problem_statement, 14, body_color)
        elif shape.shape_id == 57 and leader_name:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            format_para(p, leader_name, 16, sub_color, bold=True)

    # Let's map slides 2 to 10
    slide_shape_map = {
        2: (63, 64),
        3: (70, 71),
        4: (77, 78),
        5: (84, 85),
        6: (91, 92),
        7: (98, 99),
        8: (104, 105),
        9: (111, 112),
        10: (118, 119)
    }
    
    for s_idx, (title_id, body_id) in slide_shape_map.items():
        slide = prs.slides[s_idx - 1]
        lines = slides_data.get(s_idx, [])
        if not lines:
            continue
        
        # First line is slide title
        title_text = lines[0]
        body_lines = lines[1:]
        
        # Title Shape
        for shape in slide.shapes:
            if shape.shape_id == title_id:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                format_para(p, title_text, 28, title_color, bold=True)
                
            elif shape.shape_id == body_id:
                tf = shape.text_frame
                tf.clear()
                
                # Add paragraphs
                first = True
                for line in body_lines:
                    if not line:
                        continue
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                        
                    # Determine styling based on content
                    if line.startswith("•") or line.startswith("-"):
                        clean_line = line.lstrip("•- ").strip()
                        format_para(p, clean_line, 14, body_color, bullet=True)
                    elif re.match(r"^\d+\.", line):
                        format_para(p, line, 14, body_color)
                    elif (line.endswith(":") or 
                          "Overview" in line or 
                          "Traditional" in line or 
                          "Requirements" in line or 
                          "Evaluating" in line or 
                          "Scoring" in line or 
                          "Signals" in line or 
                          "Engine" in line or 
                          "Preventing" in line or 
                          "Suspicious" in line or 
                          "Modular" in line or 
                          "Quality" in line or 
                          "Core" in line or 
                          "Calibration" in line or 
                          "Sorter" in line or 
                          "Quantitative" in line or 
                          "Compute" in line or 
                          "Technologies" in line or 
                          "AI &" in line or 
                          "Submission" in line or 
                          "Generated" in line):
                        format_para(p, line, 16, sub_color, bold=True)
                        p.space_before = Pt(8)
                    else:
                        format_para(p, line, 14, body_color)
                        
    prs.save("presentation/presentation.pptx")
    print("Successfully generated presentation/presentation.pptx")

if __name__ == "__main__":
    main()
