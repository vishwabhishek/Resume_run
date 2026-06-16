import sys
from pptx import Presentation

pptx_path = "Idea Submission Template _ Redrob.pptx"
prs = Presentation(pptx_path)

# 1. Update Slide 5 (index 4) - Example Candidate Description
slide5 = prs.slides[4]
shape5 = next(s for s in slide5.shapes if s.name == "Google Shape;85;p17")
p3 = shape5.text_frame.paragraphs[3]
old_text5 = "".join([r.text for r in p3.runs])
print("Old Slide 5 Text:", old_text5)
new_text5 = '• Example: "Exceptional founding engineer candidate with 7.2 years experience, specializing in skills in Deep Learning, Weaviate, Recommendation Systems. Proven record building ranking/retrieval systems at Zomato, Google. locally based in Noida with a sub-30 day notice (15 days)."'
p3.text = new_text5
print("Updated Slide 5 Text:", p3.text)

# 2. Update Slide 8 (index 7) - Key Results & Shortlist Quality
slide8 = prs.slides[7]
shape8 = next(s for s in slide8.shapes if s.name == "Google Shape;105;p20")

# Update Ingestion & Filtering paragraph (index 1)
p_ingest = shape8.text_frame.paragraphs[1]
print("Old Ingest Text:", p_ingest.text)
p_ingest.text = "• Ingestion & Filtering: Ingested 100,000 candidate profiles. Hard filters rejected 9.84% of candidates, leaving 90,159 valid profiles for final ranking."
print("Updated Ingest Text:", p_ingest.text)

# Update Top Candidate 1 (index 3)
p_cand1 = shape8.text_frame.paragraphs[3]
print("Old Candidate 1 Text:", p_cand1.text)
p_cand1.text = "  - Aarav Trivedi (#1, Score: 1.1600): 7.2 yrs exp, 7.2 yrs ML exp, ex-Zomato. Skills: Deep Learning, Weaviate, Recommendation Systems. Locally based in Noida, sub-30 day notice (15 days)."
print("Updated Candidate 1 Text:", p_cand1.text)

# Update Top Candidate 2 (index 4)
p_cand2 = shape8.text_frame.paragraphs[4]
print("Old Candidate 2 Text:", p_cand2.text)
p_cand2.text = "  - Nisha Pillai (#2, Score: 1.1500): 7.6 yrs exp, 7.6 yrs ML exp, ex-Sarvam AI. Skills: scikit-learn, PyTorch, Milvus. Locally based in Gurgaon, 45-day notice."
print("Updated Candidate 2 Text:", p_cand2.text)

# Update Top Candidate 3 (index 5)
p_cand3 = shape8.text_frame.paragraphs[5]
print("Old Candidate 3 Text:", p_cand3.text)
p_cand3.text = "  - Ayaan Goyal (#3, Score: 1.1400): 6.5 yrs exp, 6.3 yrs ML exp, ex-Haptik. Skills: LoRA, Recommendation Systems, Weaviate. Locally based in Pune, sub-30 day notice (15 days)."
print("Updated Candidate 3 Text:", p_cand3.text)

# Save presentation
prs.save(pptx_path)
print("Successfully saved updated presentation!")
